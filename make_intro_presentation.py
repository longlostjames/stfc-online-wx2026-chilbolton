"""
make_intro_presentation.py
Generate an introductory PowerPoint presentation for the Chilbolton
Rainfall Climatology student project.

Usage:
    python make_intro_presentation.py [--out intro_chilbolton.pptx]
"""

from __future__ import annotations
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
NCAS_BLUE   = RGBColor(0x00, 0x3D, 0x7A)   # dark navy
NCAS_CYAN   = RGBColor(0x00, 0x9B, 0xD5)   # bright cyan
ACCENT      = RGBColor(0xFF, 0x6B, 0x00)   # orange accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF0, 0xF4, 0xF8)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p    = tf.paragraphs[0]
    p.alignment = align
    run  = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=20, color=DARK_GREY, title=None, title_size=22):
    """Add a text box with optional bold title and bulleted items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(title_size)
        run.font.bold  = True
        run.font.color.rgb = NCAS_BLUE

    for bullet in bullets:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.level = 0
        run = p.add_run()
        run.text = f'• {bullet}'
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color

    return txBox


def set_slide_background(slide, rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — title / cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NCAS_BLUE)

    # Full-width cyan banner at bottom
    add_rect(slide, 0, Inches(5.8), SLIDE_W, Inches(1.7), NCAS_CYAN)

    # Title
    add_text_box(slide, 'Chilbolton Atmospheric Observatory',
                 Inches(0.6), Inches(1.0), Inches(12.0), Inches(1.2),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_text_box(slide, 'Rainfall Climatology — Student Project',
                 Inches(0.6), Inches(2.1), Inches(12.0), Inches(0.9),
                 font_size=30, bold=False, color=NCAS_CYAN, align=PP_ALIGN.LEFT)

    # Decorative rule
    add_rect(slide, Inches(0.6), Inches(3.15), Inches(5.0), Inches(0.06), ACCENT)

    # Subtitle
    add_text_box(slide,
                 'Analysing over a decade of high-resolution\n'
                 'precipitation observations from southern England',
                 Inches(0.6), Inches(3.4), Inches(12.0), Inches(1.4),
                 font_size=20, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # Bottom banner text
    add_text_box(slide,
                 'National Centre for Atmospheric Science (NCAS)  ·  '
                 'Chilbolton, Hampshire, UK  ·  51.1°N  1.4°W',
                 Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.6),
                 font_size=14, color=NCAS_BLUE, align=PP_ALIGN.LEFT)


def slide_observatory(prs):
    """Slide 2 — What is Chilbolton Observatory?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    # Header bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NCAS_BLUE)
    add_text_box(slide, 'What is Chilbolton Atmospheric Observatory?',
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=WHITE)

    # Left column — facts
    add_bullet_box(slide,
        [
            'Located near Stockbridge, Hampshire, UK',
            'Coordinates: 51.145 °N, 1.440 °W, altitude 84 m',
            'Operated by NCAS (National Centre for Atmospheric Science)',
            'One of the UK\'s premier ground-based atmospheric observatories',
            'Continuous measurements since the early 1990s',
            'Data archived at CEDA (Centre for Environmental Data Analysis)',
            'Instruments follow NCAS/CF NetCDF conventions for interoperability',
        ],
        Inches(0.4), Inches(1.3), Inches(6.5), Inches(5.5),
        font_size=18, color=DARK_GREY,
        title='Key Facts', title_size=20)

    # Right column — context box
    add_rect(slide, Inches(7.2), Inches(1.3), Inches(5.7), Inches(5.5), WHITE)
    add_bullet_box(slide,
        [
            'Why southern England?',
            '  – Representative of maritime temperate climate',
            '  – Mix of frontal (Atlantic) and convective rainfall',
            '  – Long continuous record enables robust climatology',
            '',
            'Why high-resolution data?',
            '  – 10-second samples resolve brief intense bursts',
            '  – Enables sub-hourly intensity statistics',
            '  – Critical for flood-risk and link-budget studies',
        ],
        Inches(7.4), Inches(1.45), Inches(5.3), Inches(5.2),
        font_size=16, color=DARK_GREY)

    # Footer rule
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.06), NCAS_CYAN)


def slide_instruments(prs):
    """Slide 3 — Instruments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NCAS_BLUE)
    add_text_box(slide, 'Instruments Used in This Project',
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=WHITE)

    instruments = [
        ('Rain Gauges',
         ['Tipping-bucket gauges: 0.2 mm per tip',
          'Weighing gauges: continuous mass measurement',
          'Multiple gauges allow cross-validation',
          'QC flags identify suspect samples',
          'Labels: ncas-rain-gauge-1/2/3/5/9']),
        ('Anemometers',
         ['Sonic anemometers — no moving parts',
          'Measure 3-component wind vector at 10 s',
          'Mean wind speed and direction',
          'Used in student wind project (separate)']),
        ('Temperature & Humidity',
         ['Thermometers and capacitive RH sensors',
          'Housed in radiation screens',
          'Temperatures stored in Kelvin in NetCDF',
          'Used in student T/RH project (separate)']),
        ('Pressure',
         ['Electronic barometric pressure sensors',
          'Data available from 2019 onwards',
          'Used in student pressure project (separate)']),
    ]

    cols = 2
    col_w = Inches(6.2)
    col_h = Inches(2.6)
    pad   = Inches(0.25)

    for i, (title, bullets) in enumerate(instruments):
        col = i % cols
        row = i // cols
        lft = Inches(0.4) + col * (col_w + pad)
        top = Inches(1.3) + row * (col_h + pad)
        add_rect(slide, lft, top, col_w, col_h, WHITE)
        add_bullet_box(slide, bullets, lft + Inches(0.15), top + Inches(0.1),
                       col_w - Inches(0.3), col_h - Inches(0.2),
                       font_size=16, color=DARK_GREY,
                       title=title, title_size=18)

    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.06), NCAS_CYAN)


def slide_dataset(prs):
    """Slide 4 — The dataset."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NCAS_BLUE)
    add_text_box(slide, 'The Rainfall Dataset',
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=WHITE)

    # Stat boxes
    stats = [
        ('~12 years',     'of continuous data\n(2014 – 2026)'),
        ('10 seconds',    'sampling interval\n(6 samples min⁻¹)'),
        ('4 gauges',      'merged by physical ID\nacross sources'),
        ('NetCDF format', 'with CF conventions\nand QC flag variables'),
    ]

    box_w = Inches(2.9)
    box_h = Inches(1.8)
    gap   = Inches(0.28)
    top   = Inches(1.4)

    for i, (big, small) in enumerate(stats):
        lft = Inches(0.4) + i * (box_w + gap)
        add_rect(slide, lft, top, box_w, box_h, NCAS_BLUE)
        add_text_box(slide, big,
                     lft + Inches(0.1), top + Inches(0.15),
                     box_w - Inches(0.2), Inches(0.75),
                     font_size=26, bold=True, color=NCAS_CYAN, align=PP_ALIGN.CENTER)
        add_text_box(slide, small,
                     lft + Inches(0.1), top + Inches(0.85),
                     box_w - Inches(0.2), Inches(0.85),
                     font_size=15, color=WHITE, align=PP_ALIGN.CENTER)

    # QC section
    add_rect(slide, Inches(0.4), Inches(3.45), Inches(12.2), Inches(2.8), WHITE)
    add_bullet_box(slide,
        ['QC flag = 0 → not used (no data)',
         'QC flag = 1 → good data  ✓  (use this)',
         'QC flag = 2 → bad data (remove)',
         'QC flag ≥ 3 → bad data, specific reason (remove)',
         'Always filter to flag == 1 before computing statistics',
         'Suspect spikes may slip through — the extreme event catalogue helps identify them'],
        Inches(0.6), Inches(3.55), Inches(11.8), Inches(2.6),
        font_size=18, color=DARK_GREY,
        title='Quality Control Flags', title_size=20)

    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.06), NCAS_CYAN)


def slide_tasks_overview(prs):
    """Slide 5 — Project tasks at a glance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NCAS_BLUE)
    add_text_box(slide, 'Project Tasks — Overview',
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=WHITE)

    tasks = [
        ('1', 'Load the data',           'Read NetCDF files, apply QC flags, inspect the record'),
        ('2', 'Summary statistics',      'Total accumulation, mean/max rate, wet fraction'),
        ('3', 'Temporal aggregations',   'Daily, monthly and annual totals from raw samples'),
        ('4', 'Visualisation',           '3-panel bar chart: daily, monthly and annual'),
        ('5', 'Monthly climatology',     'Mean ± std bar chart across all years (seasonality)'),
        ('6', 'Exceedance curves',       'Probability a rate is exceeded — 1 min and 60 min'),
        ('7', 'Wet-day frequency',       'P(rain on a given day) — overall and by month'),
        ('8', 'Rainfall intensity',      'Box plots and percentile curves by month; per-year heatmap'),
    ]

    row_h   = Inches(0.68)
    top_off = Inches(1.25)
    num_w   = Inches(0.55)
    title_w = Inches(3.2)
    desc_w  = Inches(8.8)
    lft     = Inches(0.4)

    for i, (num, title, desc) in enumerate(tasks):
        top = top_off + i * row_h
        bg  = WHITE if i % 2 == 0 else LIGHT_GREY
        add_rect(slide, lft, top, Inches(12.5), row_h - Inches(0.04), bg)

        # Number bubble
        add_rect(slide, lft + Inches(0.07), top + Inches(0.08),
                 Inches(0.44), Inches(0.5), NCAS_CYAN)
        add_text_box(slide, num,
                     lft + Inches(0.07), top + Inches(0.08),
                     Inches(0.44), Inches(0.5),
                     font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     lft + num_w + Inches(0.1), top + Inches(0.1),
                     title_w, Inches(0.5),
                     font_size=17, bold=True, color=NCAS_BLUE)

        add_text_box(slide, desc,
                     lft + num_w + title_w + Inches(0.1), top + Inches(0.1),
                     desc_w, Inches(0.5),
                     font_size=16, color=DARK_GREY)

    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.06), NCAS_CYAN)


def slide_learning_outcomes(prs):
    """Slide 6 — Learning outcomes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, LIGHT_GREY)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NCAS_BLUE)
    add_text_box(slide, 'What You Will Learn',
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=WHITE)

    cols = [
        ('Data Skills',
         ['Read and explore NetCDF files with Python',
          'Apply quality control to real observations',
          'Aggregate 10-second data to daily/monthly/annual timescales',
          'Build an event catalogue with de-duplication logic']),
        ('Statistical Methods',
         ['Compute descriptive statistics (mean, std, percentiles)',
          'Construct empirical exceedance (CCDF) curves',
          'Interpret high-percentile statistics (99th, 99.9th)',
          'Analyse inter-annual variability with heatmaps']),
        ('Physical Understanding',
         ['Distinguish rainfall accumulation from intensity',
          'Identify the seasonal cycle in frequency vs. intensity',
          'Recognise convective vs. frontal rainfall regimes',
          'Link gauge data to named storms and weather events']),
        ('Scientific Communication',
         ['Produce publication-quality multi-panel figures',
          'Label axes with correct units (mm hr⁻¹, mm day⁻¹)',
          'Interpret and discuss results in the project report',
          'Cross-reference data with independent sources']),
    ]

    col_w = Inches(6.0)
    col_h = Inches(2.7)
    pad   = Inches(0.25)

    for i, (title, bullets) in enumerate(cols):
        col = i % 2
        row = i // 2
        lft = Inches(0.4) + col * (col_w + pad)
        top = Inches(1.3) + row * (col_h + pad)
        add_rect(slide, lft, top, col_w, col_h, WHITE)
        # Coloured top strip
        add_rect(slide, lft, top, col_w, Inches(0.12), NCAS_CYAN)
        add_bullet_box(slide, bullets,
                       lft + Inches(0.15), top + Inches(0.18),
                       col_w - Inches(0.3), col_h - Inches(0.3),
                       font_size=16, color=DARK_GREY,
                       title=title, title_size=18)

    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.06), NCAS_CYAN)


def slide_getting_started(prs):
    """Slide 7 — Getting started."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, NCAS_BLUE)

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NCAS_CYAN)
    add_text_box(slide, 'Getting Started',
                 Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.85),
                 font_size=28, bold=True, color=NCAS_BLUE)

    steps = [
        ('1', 'Open the introduction notebooks',
               'intro_chilbolton.ipynb  →  observatory overview and dataset summary\n'
               'intro_netcdf.ipynb     →  how to read NetCDF files with Python\n'
               'intro_statistics.ipynb →  statistical methods used in the project'),
        ('2', 'Open your worksheet',
               'student1_rainfall.ipynb\n'
               'Work through Tasks 1 – 8 in order. Each task builds on the previous one.'),
        ('3', 'Run the setup cell first',
               'Cell 3 defines all helper functions and the GaugeData dataclass.\n'
               'Run it once before any other cell.'),
        ('4', 'Ask for help',
               'Refer to the introduction notebooks, Python/NumPy/Pandas documentation,\n'
               'or ask your supervisor if you get stuck.'),
    ]

    top_off = Inches(1.35)
    row_h   = Inches(1.38)

    for i, (num, title, body) in enumerate(steps):
        top = top_off + i * row_h
        add_rect(slide, Inches(0.4), top, Inches(12.2), row_h - Inches(0.08), RGBColor(0x00, 0x2A, 0x55))

        # Number
        add_rect(slide, Inches(0.5), top + Inches(0.18), Inches(0.55), Inches(0.55), ACCENT)
        add_text_box(slide, num,
                     Inches(0.5), top + Inches(0.18), Inches(0.55), Inches(0.55),
                     font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text_box(slide, title,
                     Inches(1.2), top + Inches(0.12), Inches(11.2), Inches(0.42),
                     font_size=19, bold=True, color=NCAS_CYAN)
        add_text_box(slide, body,
                     Inches(1.2), top + Inches(0.55), Inches(11.2), Inches(0.75),
                     font_size=15, color=LIGHT_GREY)

    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.06), ACCENT)


# ── Main ──────────────────────────────────────────────────────────────────────

def build(out_path: Path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_observatory(prs)
    slide_instruments(prs)
    slide_dataset(prs)
    slide_tasks_overview(prs)
    slide_learning_outcomes(prs)
    slide_getting_started(prs)

    prs.save(str(out_path))
    print(f'Saved: {out_path}  ({out_path.stat().st_size // 1024} KB)')


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate Chilbolton intro presentation')
    parser.add_argument('--out', default='intro_chilbolton.pptx',
                        help='Output .pptx filename (default: intro_chilbolton.pptx)')
    args = parser.parse_args()
    build(Path(args.out))
